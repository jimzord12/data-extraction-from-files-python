from dotenv import load_dotenv
from pptx import Presentation
import os

# Load the .env file
load_dotenv()

pptx_dir_path = os.getenv("PPTX_DIR_PATH")
output_file = os.getenv("OUTPUT_FILE_PATH")

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_content = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)

    return "\n".join(text_content)


def extract_text_from_directory():
    combined_text = []

    # Walk through the directory and find all .pptx files
    for root, _, files in os.walk(pptx_dir_path):
        for file in files:
            if file.endswith(".pptx"):
                file_path = os.path.join(root, file)
                print(f"Extracting text from: {file_path}")
                extracted_text = extract_text_from_pptx(file_path)
                
                # Add a header for each file for clarity
                combined_text.append(f"\n--- Extracted from {file} ---\n")
                combined_text.append(extracted_text)

    # Combine all extracted text and save it to the output file
    all_text = "\n".join(combined_text)
    save_text_to_file(all_text)

def save_text_to_file(text):
    with open(output_file, "w", encoding="utf-8") as f:
        f.write(text)
    print(f"Text content saved to {output_file}")


if (__name__ == "__main__"):
    print("Extracting text from .pptx files...")
    print("The directory path is: ", pptx_dir_path)
    print("")
    extract_text_from_directory()

# # Example usage
# file_path = "example.pptx"
# text = extract_text_from_pptx(file_path)

# # Create a .txt file with the extracted text
# output_file = "extracted_content.txt"
# save_text_to_file(text, output_file)